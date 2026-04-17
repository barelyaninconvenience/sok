"""Extract all text content from a pptx, slide by slide, with notes."""
import sys
from pptx import Presentation

path = sys.argv[1]
prs = Presentation(path)

print(f"=== {path} ===")
print(f"Slides: {len(prs.slides)}\n")

for i, slide in enumerate(prs.slides, 1):
    print(f"\n{'=' * 70}")
    print(f"SLIDE {i}")
    print('=' * 70)

    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                print(text)
                print()

    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            print(f"--- SPEAKER NOTES ---")
            print(notes)
            print()
