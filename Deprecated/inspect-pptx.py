"""Inspect a .pptx file to verify slide structure and content."""
import sys
from pptx import Presentation

path = sys.argv[1]
prs = Presentation(path)

print(f"File: {path}")
print(f"Slide count: {len(prs.slides)}")
print(f"Slide dimensions: {prs.slide_width} x {prs.slide_height} EMU")
print()
print("=== Slide titles + shape counts ===")
for i, slide in enumerate(prs.slides, 1):
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text and title is None:
                title = text.split("\n")[0][:100]
                break
    shape_count = len(slide.shapes)
    text_total = sum(1 for s in slide.shapes if s.has_text_frame)
    print(f"Slide {i:2d}: [{shape_count} shapes, {text_total} text frames] {title or '(no title)'}")

print()
print("=== First slide text sample ===")
first = prs.slides[0]
for shape in first.shapes:
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                print(f"  | {t[:120]}")
