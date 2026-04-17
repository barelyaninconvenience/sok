"""Extract all text content from a pptx to a markdown file (bypasses console encoding)."""
import sys
from pathlib import Path
from pptx import Presentation

path = Path(sys.argv[1])
prs = Presentation(str(path))
out_path = path.with_suffix(".extracted.md")

lines = []
lines.append(f"# Extracted slides: {path.name}")
lines.append(f"")
lines.append(f"Total slides: {len(prs.slides)}")
lines.append(f"")

for i, slide in enumerate(prs.slides, 1):
    lines.append(f"")
    lines.append(f"---")
    lines.append(f"")
    lines.append(f"## SLIDE {i}")
    lines.append(f"")
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if text:
                lines.append(text)
                lines.append("")
    if slide.has_notes_slide:
        notes = slide.notes_slide.notes_text_frame.text.strip()
        if notes:
            lines.append(f"**SPEAKER NOTES:**")
            lines.append(f"")
            lines.append(notes)
            lines.append("")

out_path.write_text("\n".join(lines), encoding="utf-8")
print(f"Wrote {out_path} ({out_path.stat().st_size} bytes, {len(prs.slides)} slides)")
