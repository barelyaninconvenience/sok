"""Prep IT7021C submission .pptx from the markdown outline.
Steps:
  1. Read the original .md
  2. Strip the instruction blockquote + surrounding horizontal rules
  3. Replace em-dashes and smart quotes with ASCII equivalents
  4. Write to _submission.md
  5. Run pandoc
  6. Inspect result
"""
import re
import subprocess
from pathlib import Path

src_dir = Path(r"C:\Users\shelc\Documents\UC MS-IS\SPRING 2026\IT7021C ENTERPRISE SECURITY FORENSICS")
src_md = src_dir / "Caddell_IT7021C_ProjectPresentation.md"
sub_md = src_dir / "Caddell_IT7021C_ProjectPresentation_submission.md"
out_pptx = src_dir / "Caddell_IT7021C_ProjectPresentation_generated.pptx"
pandoc = r"C:\Program Files\Pandoc\pandoc.exe"

text = src_md.read_text(encoding="utf-8")

# Strip the instruction block + the --- separators around it
# The block is:
#   ---
#
#   > **PRESENTATION OUTLINE ...**
#   > Target: 10-15 minutes...
#
#   ---
pattern = r"---\s*\n\s*\n> \*\*PRESENTATION OUTLINE.*?\n> Target:.*?\n\s*\n---\s*\n"
text = re.sub(pattern, "", text, flags=re.DOTALL)

# Encoding normalization for pandoc pptx writer compatibility
replacements = {
    "\u2014": "-",  # em dash
    "\u2013": "-",  # en dash
    "\u2018": "'",  # left single quote
    "\u2019": "'",  # right single quote
    "\u201c": '"',  # left double quote
    "\u201d": '"',  # right double quote
    "\u2026": "...",  # ellipsis
}
for k, v in replacements.items():
    text = text.replace(k, v)

sub_md.write_text(text, encoding="utf-8")
print(f"Wrote submission md: {sub_md.name} ({len(text)} chars, {text.count(chr(10))} lines)")

# Confirm instruction block is gone
if "PRESENTATION OUTLINE" in text:
    print("WARNING: instruction block still present in submission md")
else:
    print("OK: instruction block removed")

# Run pandoc
result = subprocess.run(
    [pandoc, "-f", "markdown", "-t", "pptx", str(sub_md), "-o", str(out_pptx)],
    capture_output=True,
    text=True,
)
if result.returncode != 0:
    print(f"Pandoc failed: {result.stderr}")
    raise SystemExit(1)
print(f"Pandoc OK -> {out_pptx.name} ({out_pptx.stat().st_size} bytes)")

# Inspect the result
from pptx import Presentation
prs = Presentation(str(out_pptx))
print(f"\n=== {len(prs.slides)} slides ===")
for i, slide in enumerate(prs.slides, 1):
    title = None
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()
            if t and title is None:
                title = t.split("\n")[0][:80]
                break
    print(f"  {i:2d}: {title or '(no title)'}")
